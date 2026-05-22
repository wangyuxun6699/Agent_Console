"""文档加载和分片服务 - 支持多格式与多模态图片识别"""
import os
import uuid
import pathlib
import fitz  # PyMuPDF
import pptx
import docx
import pandas as pd
from typing import Dict, List

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.chat_models.tongyi import ChatTongyi
from langchain_core.messages import HumanMessage


class DocumentLoader:
    """文档加载和分片服务（支持 PDF/PPTX/Word/Excel/TXT 及图片 OCR 描述）"""

    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50, image_output_dir: str = "./extracted_images"):
        # 1. 初始化三层滑动窗口分块器
        level_1_size = max(1200, chunk_size * 2)
        level_1_overlap = max(240, chunk_overlap * 2)
        level_2_size = max(600, chunk_size)
        level_2_overlap = max(120, chunk_overlap)
        level_3_size = max(300, chunk_size // 2)
        level_3_overlap = max(60, chunk_overlap // 2)

        separators = ["\n\n", "\n", "。", "！", "？", "，", "、", " ", ""]
        self._splitter_level_1 = RecursiveCharacterTextSplitter(chunk_size=level_1_size, chunk_overlap=level_1_overlap, add_start_index=True, separators=separators)
        self._splitter_level_2 = RecursiveCharacterTextSplitter(chunk_size=level_2_size, chunk_overlap=level_2_overlap, add_start_index=True, separators=separators)
        self._splitter_level_3 = RecursiveCharacterTextSplitter(chunk_size=level_3_size, chunk_overlap=level_3_overlap, add_start_index=True, separators=separators)

        # 2. 初始化图片临时存储目录和视觉大模型
        self.image_output_dir = image_output_dir
        if not os.path.exists(self.image_output_dir):
            os.makedirs(self.image_output_dir)
            
        self.vlm_image = ChatTongyi(model="qwen-vl-plus-2025-05-07", max_retries=3)

    def _clean_vlm_output(self, content) -> str:
        """脱壳逻辑：将大模型返回的列表转为纯字符串"""
        if isinstance(content, list):
            return "".join([item.get("text", "") for item in content if isinstance(item, dict)]).strip()
        return str(content).strip()

    def _image_to_text_summary(self, img_path: str) -> str:
        """调用多模态模型识别图片/扫描件内容"""
        try:
            abs_path = str(pathlib.Path(img_path).absolute())
            message = HumanMessage(content=[
                {"text": "请详细描述这张图片的内容，包括图表数据、文字或核心信息，以便录入知识库被检索。"},
                {"image": f"file://{abs_path}"}
            ])
            print(f"--- 正在调用 Qwen-VL-PLUS 分析图片: {os.path.basename(img_path)} ---")
            response = self.vlm_image.invoke([message])
            return self._clean_vlm_output(response.content)
        except Exception as e:
            print(f"图片识别失败 {img_path}: {e}")
            return ""

    @staticmethod
    def _build_chunk_id(filename: str, page_number: int, level: int, index: int) -> str:
        return f"{filename}::p{page_number}::l{level}::{index}"

    def _split_page_to_three_levels(self, text: str, base_doc: Dict, page_global_chunk_idx: int) -> List[Dict]:
        """将提取出的单页文本（含图片描述）进行三层分块"""
        if not text:
            return []

        root_chunks: List[Dict] = []
        page_number = int(base_doc.get("page_number", 0))
        filename = base_doc["filename"]

        level_1_docs = self._splitter_level_1.create_documents([text], [base_doc])
        level_1_counter, level_2_counter, level_3_counter = 0, 0, 0

        for level_1_doc in level_1_docs:
            level_1_text = (level_1_doc.page_content or "").strip()
            if not level_1_text:
                continue
            level_1_id = self._build_chunk_id(filename, page_number, 1, level_1_counter)
            level_1_counter += 1

            level_1_chunk = {**base_doc, "text": level_1_text, "chunk_id": level_1_id, "parent_chunk_id": "", "root_chunk_id": level_1_id, "chunk_level": 1, "chunk_idx": page_global_chunk_idx}
            page_global_chunk_idx += 1
            root_chunks.append(level_1_chunk)

            level_2_docs = self._splitter_level_2.create_documents([level_1_text], [base_doc])
            for level_2_doc in level_2_docs:
                level_2_text = (level_2_doc.page_content or "").strip()
                if not level_2_text:
                    continue
                level_2_id = self._build_chunk_id(filename, page_number, 2, level_2_counter)
                level_2_counter += 1

                level_2_chunk = {**base_doc, "text": level_2_text, "chunk_id": level_2_id, "parent_chunk_id": level_1_id, "root_chunk_id": level_1_id, "chunk_level": 2, "chunk_idx": page_global_chunk_idx}
                page_global_chunk_idx += 1
                root_chunks.append(level_2_chunk)

                level_3_docs = self._splitter_level_3.create_documents([level_2_text], [base_doc])
                for level_3_doc in level_3_docs:
                    level_3_text = (level_3_doc.page_content or "").strip()
                    if not level_3_text:
                        continue
                    level_3_id = self._build_chunk_id(filename, page_number, 3, level_3_counter)
                    level_3_counter += 1
                    root_chunks.append({**base_doc, "text": level_3_text, "chunk_id": level_3_id, "parent_chunk_id": level_2_id, "root_chunk_id": level_1_id, "chunk_level": 3, "chunk_idx": page_global_chunk_idx})
                    page_global_chunk_idx += 1

        return root_chunks

    def load_document(self, file_path: str, filename: str) -> list[dict]:
        """按文件类型智能解析并提取特征图文组合分片"""
        file_lower = filename.lower()
        documents = []
        page_global_chunk_idx = 0
        
        try:
            # 1. TXT 处理
            if file_lower.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as f:
                    text_content = f.read()
                base_doc = {"filename": filename, "file_path": file_path, "file_type": "TXT", "page_number": 1}
                chunks = self._split_page_to_three_levels(text_content, base_doc, page_global_chunk_idx)
                documents.extend(chunks)

            # 2. PDF 处理（含提取图片与 OCR 扫面件识别）
            elif file_lower.endswith(".pdf"):
                doc = fitz.open(file_path)
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    page_texts = [page.get_text()]
                    
                    # 提取内嵌图片进行模型识别
                    for img in page.get_images(full=True):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        ext = base_image["ext"]
                        img_filename = f"{uuid.uuid4().hex}.{ext}"
                        img_filepath = os.path.join(self.image_output_dir, img_filename)
                        
                        with open(img_filepath, "wb") as f:
                            f.write(base_image["image"])
                            
                        img_desc = self._image_to_text_summary(img_filepath)
                        if img_desc:
                            page_texts.append(f"\n[第{page_num+1}页图片描述]: {img_desc}\n")
                            
                    full_page_text = "\n".join(page_texts).strip()
                    base_doc = {"filename": filename, "file_path": file_path, "file_type": "PDF", "page_number": page_num + 1}
                    chunks = self._split_page_to_three_levels(full_page_text, base_doc, page_global_chunk_idx)
                    page_global_chunk_idx += len(chunks)
                    documents.extend(chunks)

            # 3. PPTX 处理（含文字和形状内的图片）
            elif file_lower.endswith((".pptx", ".ppt")):
                prs = pptx.Presentation(file_path)
                for i, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_texts.append(shape.text)
                        if hasattr(shape, "image"):
                            ext = shape.image.ext
                            img_filename = f"{uuid.uuid4().hex}.{ext}"
                            img_filepath = os.path.join(self.image_output_dir, img_filename)
                            with open(img_filepath, "wb") as f:
                                f.write(shape.image.blob)
                            img_desc = self._image_to_text_summary(img_filepath)
                            if img_desc:
                                slide_texts.append(f"\n[第{i+1}页幻灯片图片描述]: {img_desc}\n")
                                
                    full_slide_text = "\n".join(slide_texts).strip()
                    base_doc = {"filename": filename, "file_path": file_path, "file_type": "PPT", "page_number": i + 1}
                    chunks = self._split_page_to_three_levels(full_slide_text, base_doc, page_global_chunk_idx)
                    page_global_chunk_idx += len(chunks)
                    documents.extend(chunks)

            # 4. DOCX 处理
            elif file_lower.endswith((".docx", ".doc")):
                doc = docx.Document(file_path)
                full_text = "\n".join([para.text for para in doc.paragraphs])
                base_doc = {"filename": filename, "file_path": file_path, "file_type": "Word", "page_number": 1}
                chunks = self._split_page_to_three_levels(full_text, base_doc, page_global_chunk_idx)
                documents.extend(chunks)

            # 5. Excel / CSV 处理
            elif file_lower.endswith((".xlsx", ".xls", ".csv")):
                df = pd.read_excel(file_path) if file_lower.endswith(('.xlsx', '.xls')) else pd.read_csv(file_path)
                full_text = df.to_string()
                base_doc = {"filename": filename, "file_path": file_path, "file_type": "Excel", "page_number": 1}
                chunks = self._split_page_to_three_levels(full_text, base_doc, page_global_chunk_idx)
                documents.extend(chunks)

            else:
                raise ValueError(f"不支持的文件类型: {filename}")

            return documents
        except Exception as e:
            raise Exception(f"处理文档失败 ({filename}): {str(e)}")

    def load_documents_from_folder(self, folder_path: str) -> list[dict]:
        """批量遍历处理目录结构"""
        all_documents = []
        valid_exts = (".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".csv", ".txt")
        
        for filename in os.listdir(folder_path):
            file_lower = filename.lower()
            if not file_lower.endswith(valid_exts):
                continue
            file_path = os.path.join(folder_path, filename)
            try:
                print(f"开始加载文档: {filename}")
                documents = self.load_document(file_path, filename)
                all_documents.extend(documents)
            except Exception as e:
                print(f"跳过文件 {filename}, 错误: {e}")
                continue
                
        return all_documents